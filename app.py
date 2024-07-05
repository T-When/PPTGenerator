import streamlit as st
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
import vertexai
from vertexai.generative_models import GenerativeModel, Part, FinishReason
import vertexai.preview.generative_models as generative_models
from vertexai.preview.vision_models import ImageGenerationModel
import json
import os


def format_list_with_textbox_style(presentation, slide_num, text_box_id, list_items):
    slide = presentation.slides[slide_num - 1]
    for shape in slide.shapes:
        if shape.name == text_box_id:
            text_frame = shape.text_frame
            # Clear existing text
            text_frame.clear()

            # Add paragraphs with preserved formatting
            for item in list_items:
                p = text_frame.add_paragraph()
                p.text = f"• {item}"

                # Preserve formatting of the first run
                if text_frame.paragraphs[0].runs:
                    first_run = text_frame.paragraphs[0].runs[0]
                    font = first_run.font
                    p.font.name = font.name
                    p.font.size = font.size
                    p.font.bold = font.bold
                    p.font.italic = font.italic
                    p.font.underline = font.underline
                    if font.color.rgb is not None:
                        p.font.color.rgb = font.color.rgb

                # Preserve paragraph alignment
                p.alignment = shape.text_frame.paragraphs[0].alignment

                # Increase line spacing
                p.space_before = Pt(1)
                p.space_after = Pt(0)
                p.line_spacing = Pt(24)  # Adjust as needed (e.g., 24 points)
    return presentation

def update_text_of_textbox(presentation, slide_number, text_box_id, new_text):
    if presentation is None:
        raise ValueError("Presentation is None. Please check if the PowerPoint file was loaded correctly.")

    # Adjust slide number (1-based index in PowerPoint slides)
    slide_index = slide_number - 1

    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise ValueError(f"Slide number {slide_number} is out of range.")

    slide = presentation.slides[slide_index]

    text_box_found = False

    for shape in slide.shapes:
        if shape.name == text_box_id:
            text_box_found = True
            text_frame = shape.text_frame
            # Preserve formatting of the first run
            first_paragraph = text_frame.paragraphs[0]
            first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
            font = first_run.font
            font_name = font.name
            font_size = font.size
            font_bold = font.bold
            font_italic = font.italic
            font_underline = font.underline
            
            try:
                font_color = font.color.rgb
            except AttributeError:
                font_color = None
            
            text_frame.clear() 
            new_run = text_frame.paragraphs[0].add_run() 
            new_run.text = new_text
            new_run.font.name = font_name
            new_run.font.size = font_size
            new_run.font.bold = font_bold
            new_run.font.italic = font_italic
            new_run.font.underline = font_underline
            if font_color:
                new_run.font.color.rgb = font_color
        else :
            presentation = presentation

    

    return presentation



def replace_image_in_ppt(presentation, slide_number, picturebox_id, new_image_path):
    
    
    slide = presentation.slides[slide_number - 1]
    
    
    picture_shape = None
    for shape in slide.shapes:
        if shape.name ==  picturebox_id:
            picture_shape = shape
            break
    
    if picture_shape is None:
        raise ValueError(f"No picture shape found with ID {picturebox_id} on slide {slide_number}")
    
    
    left = picture_shape.left
    top = picture_shape.top
    width = picture_shape.width
    height = picture_shape.height
    
    # Remove the old picture
    sp = picture_shape._element
    sp.getparent().remove(sp)
    
    
    slide.shapes.add_picture(new_image_path, left, top, width, height)
    
    return presentation
    



def img_generation(prompt):

    project_id = "corp-tas-poc-ai"
    output_file = "output.png"


    vertexai.init(project=project_id, location="us-central1")

    model = ImageGenerationModel.from_pretrained("imagegeneration@006")

    images = model.generate_images(
        prompt=prompt,
        number_of_images=1,
        language="en",
        aspect_ratio="3:4",
        safety_filter_level="block_some",
        person_generation="allow_adult",
    )

    images[0].save(location=output_file, include_generation_parameters=False)


    print(f"Created output image using {len(images[0]._image_bytes)} bytes")

    return True






def generate(input_prompt):
  
    generation_config = {
    "max_output_tokens": 8192,
    "temperature": 1,
    "top_p": 0.95,
}

    safety_settings = {
    generative_models.HarmCategory.HARM_CATEGORY_HATE_SPEECH: generative_models.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
    generative_models.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: generative_models.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
    generative_models.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: generative_models.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
    generative_models.HarmCategory.HARM_CATEGORY_HARASSMENT: generative_models.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
}
  
    vertexai.init(project="corp-tas-poc-ai", location="us-central1")
    model = GenerativeModel(
    "gemini-1.5-flash-001",
    )
    responses = model.generate_content(
        [input_prompt],
        generation_config=generation_config,
        safety_settings=safety_settings,
        stream=True,
    )
    combined_text =''
    for response in responses:
        #print(response.text, end="")
        combined_text += response.text
    return combined_text


def prompt(context):
    input_prompt = f"""you will be providing all the presentation details of the given topic: {context}

follow these instructions:
slide 1 : Title Slide : Title, Subtitle
slide 2 Agenda Slide : Title, Subtitle, Bullet Points
slide 3 and 4 Content Slides  : Headings, Bullet Points,image
slide 5 Conclusion Slide :  Title, Subtitle, Bullet Points
slide 6 Thank You Slide : Title, Text

Output should be in JSON structure
    
{{
    "slides": [
        {{
            "slide_number": 1,
            "type": "Title Slide",
            "title": "",
            "subtitle": ""
        }},
        {{
            "slide_number": 2,
            "type": "Agenda Slide",
            "title": "",
            "bullet_points": ["", ""]
        }},
        {{
            "slide_number": 3,
            "type": "Content Slide",
            "title": "",
            "heading": "",
            "bullet_points": ["", ""]
            "image":"describe the image detailedly"
        }},
        {{
            "slide_number": 4,
            "type": "Content Slide",
            "title": "",
            "heading": "",
            "bullet_points": ["", ""]
            "image":"describe the image detailedly"

        }},
        {{
            "slide_number": 5,
            "type": "Conclusion Slide",
            "title": "",
            "subtitle": "",
            "bullet_points": ["", ""]
        }},
        {{
            "slide_number": 6,
            "type": "Thank You Slide",
            "title": "",
            "text": ""
        }}
    ]
}}"""
    return input_prompt



def main():
    st.title("PPT Generator")
    context = st.text_area("Enter some text:", height=200)
    if st.button("Generate PowerPoint"):
        input_prompt_2 = prompt(context)
        
        out = generate(input_prompt_2)
        print(out)
        json_string = out[7:-3]
        print(json_string)
        data = json.loads(json_string)
        presentation = pptx.Presentation(os.path.join("powerpoints", "genpact_temp_2.pptx"))
        
        if presentation is None:
            raise ValueError("Failed to load PowerPoint presentation.")
        try:
            for i in range(len(data["slides"])):
                for key in data["slides"][i].keys():
                    if key == "slide_number" or key == "type":
                        continue
                    elif key.lower() == "bullet_points":
                        slide_number = i + 1
                        text_box_id = key.lower()
                        my_lists = data["slides"][i][key]
                        print(slide_number, text_box_id, my_lists)
                        st.write(text_box_id)
                        st.write(my_lists)

                        format_list_with_textbox_style(presentation, slide_number, text_box_id, my_lists)

                    elif key.lower() == "image":
                        slide_number = i + 1
                        picturebox_id =key.lower()
                        img_prompt = data["slides"][i][key]
                        image = img_generation(img_prompt)
                        new_image_path = "output.png"

                        if image:
                            replace_image_in_ppt(presentation, slide_number, picturebox_id, new_image_path)
                        

                        

                    
                    else:
                        slide_number = i + 1
                        text_box_id = key.lower()
                        new_text = data["slides"][i][key]
                        print(slide_number, text_box_id, new_text)
                        st.write(text_box_id)
                        st.write(new_text)
                        presentation = update_text_of_textbox(presentation, slide_number, text_box_id, new_text)

            output_path = "output.pptx"
            
            presentation.save("output.pptx")
            # st.markdown(f"Download your PowerPoint presentation [here](/{output_path}).")
            # st.download_button(label="Download PowerPoint", data=output_path, file_name="output.pptx")
            try :
                with open(output_path, 'rb') as f:
                    pptx_bytes = f.read()
                st.download_button(label="Download PowerPoint", data=pptx_bytes, file_name="output.pptx")
            except Exception as e:
                print(f"An error occurred: {str(e)}")


            print("Completed")
        
        except Exception as e:
            print(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()

    
